"""Xuất bộ slide ra file .pptx sửa được bằng PowerPoint / LibreOffice / Google Slides.

Vì sao có cả bản này lẫn bản HTML: bản HTML giữ đúng thiết kế và in ra PDF được,
nhưng không sửa được. Nhiều hội đồng và hội nghị trong nước lại đòi nộp .pptx, và
người trình bày thường muốn tự chỉnh vài chữ ngay trước khi lên nói.

Ba chỗ python-pptx KHÔNG làm được, và cách xử lý ở đây:

- **Không có autofit thật.** `TextFrame.fit_text()` phụ thuộc vào việc đo font
  ngoài thư viện và hay tràn. Nên mọi khung chữ ở đây đều được đặt toạ độ và cỡ
  chữ bằng tay, cùng bộ số với `_SLIDES_CSS` (1px trên khung 1280×720 = 0,75pt).
- **Không render được Mermaid.** Nhưng `DIAGRAM_RULES` đã giới hạn sơ đồ ở mức
  `flowchart TD|LR`, ≤9 node, nhãn ≤8 chữ — đủ đơn giản để **vẽ lại bằng shape
  gốc của PowerPoint**. Làm vậy còn hơn nhúng ảnh: người dùng kéo, sửa chữ, đổi
  màu được. Xem `_draw_diagram`.
- **Không có công thức.** Công thức ở đây lưu dạng `^{…}` / `_{…}` nên dựng được
  bằng chỉ số trên/dưới ở mức run (`_rich_runs`), không cần OMML.
"""

from __future__ import annotations

import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from . import store
from . import slide_theme as theme

# Khung 16:9 giống hệt bản HTML: 1280×720 px @96dpi = 13,333in × 7,5in.
# 720px = 540pt nên 1px = 0,75pt — mọi con số dưới đây quy từ `_SLIDES_CSS`.
PX = 0.75
W_PX, H_PX = 1280, 720
PAD_X, PAD_TOP, PAD_BOT = 60, 44, 56

INK = RGBColor(0x0F, 0x17, 0x2A)
INK_2 = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
LINE = RGBColor(0xDF, 0xE1, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"        # có mặt trên mọi máy, dấu tiếng Việt đầy đủ
MONO = "Consolas"

# Dấu tiếng Việt chồng tầng (ế, ộ, ữ) bị cắt ngọn dưới 1.28 — xem CLAUDE.md
LINE_SPACING = 1.28

_SUP = re.compile(r"\^\{([^{}]*)\}")
_SUB = re.compile(r"_\{([^{}]*)\}")
_TOKEN = re.compile(r"\^\{[^{}]*\}|_\{[^{}]*\}")


def px(v: float) -> Pt:
    return Pt(v * PX)


# ------------------------------------------------------------------ chữ


def _rich_runs(para, text: str, *, size: float, color=INK, bold=False,
               font=FONT) -> None:
    """Đổ chữ vào một paragraph, dựng `^{…}` / `_{…}` thành chỉ số thật.

    Đây là bản anh em của `sci()` bên app.js và `rich()` bên main.py: dạng
    `^{…}` là để model đọc, không phải để người nhìn.
    """
    def add(chunk: str, baseline: int = 0) -> None:
        if not chunk:
            return
        r = para.add_run()
        r.text = chunk
        r.font.size = px(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
        if baseline:
            # python-pptx chưa bọc superscript/subscript — đặt thẳng thuộc tính
            # `baseline` trên rPr, đúng như OOXML quy định (đơn vị 1/1000 %)
            r.font._rPr.set("baseline", "30000" if baseline > 0 else "-25000")

    pos = 0
    for m in _TOKEN.finditer(text or ""):
        add(text[pos:m.start()])
        tok = m.group(0)
        if tok.startswith("^"):
            add(_SUP.match(tok).group(1), 1)
        else:
            add(_SUB.match(tok).group(1), -1)
        pos = m.end()
    add((text or "")[pos:])


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _text(slide, x, y, w, h, text, *, size, color=INK, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=LINE_SPACING):
    tf = _box(slide, x, y, w, h)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    _rich_runs(p, text, size=size, color=color, bold=bold, font=font)
    return tf


def _bullets(slide, x, y, w, h, items, *, size, color=INK, gap=10,
             marker="•", anchor=MSO_ANCHOR.TOP):
    """Gạch đầu dòng. Tự vẽ dấu đầu dòng thay vì dùng danh sách của PowerPoint —
    danh sách thật kéo theo cả bộ thụt lề của theme, mà ta không dùng theme."""
    tf = _box(slide, x, y, w, h)
    tf.vertical_anchor = anchor
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = LINE_SPACING
        p.space_after = px(gap)
        _rich_runs(p, (f"{marker}  " if marker else "") + it, size=size, color=color)
    return tf


# -------------------------------------------------------------- sơ đồ


# Khai node: `A["nhãn"]`, `C{"câu hỏi?"}`, `B(nhãn)`. Ngoặc phải dính liền mã
# node, không có khoảng trắng — đúng cú pháp mà `DIAGRAM_RULES` bắt buộc.
_MMD_DECL = re.compile(r'\b([A-Za-z][A-Za-z0-9]*)([\[({]{1,2})\s*"?(.*?)"?\s*([\])}]{1,2})')
# Nhãn trên cạnh: `-->|"ghi chú"|`. Phải bóc RA TRƯỚC, nếu không nó bị đọc nhầm
# thành khai node và sinh ra mấy node rác tên "u", "ch".
_MMD_ELABEL = re.compile(r'\|\s*"?([^"|]*)"?\s*\|')
_MMD_ARROW = re.compile(r'-{2,3}>|==+>|-\.->')


def parse_mermaid(code: str) -> tuple[str, dict[str, str], list[tuple[str, str, str]]]:
    """Bóc `flowchart` thành (hướng, {mã node: nhãn}, [(từ, đến, nhãn cạnh)]).

    Chỉ hiểu đúng tập con mà `DIAGRAM_RULES` cho phép — cố tình không viết bộ
    phân tích Mermaid đầy đủ, vì mọi sơ đồ ở đây đều do prompt đó sinh ra.
    """
    lines = [l for l in (code or "").splitlines() if l.strip()]
    if not lines:
        return "TD", {}, []
    direction = "LR" if lines[0].strip().upper().endswith("LR") else "TD"
    nodes: dict[str, str] = {}
    edges: list[tuple[str, str, str]] = []

    for raw in lines[1:]:
        line = raw.strip()
        if not line or line.startswith("%%"):
            continue

        # 1. bóc nhãn cạnh ra khỏi dòng trước đã
        elabel = ""
        if (m := _MMD_ELABEL.search(line)):
            elabel = (m.group(1) or "").strip()
            line = line[:m.start()] + " " + line[m.end():]

        # 2. ghi nhận mọi khai node trong dòng, rồi thay bằng chính mã node
        def take(m: re.Match) -> str:
            nid, label = m.group(1), (m.group(3) or "").strip()
            if label:
                nodes[nid] = label
            else:
                nodes.setdefault(nid, nid)
            return " " + nid + " "

        line = _MMD_DECL.sub(take, line)

        # 3. phần còn lại chỉ còn mã node và mũi tên
        if _MMD_ARROW.search(line):
            chain = [p.strip() for p in _MMD_ARROW.split(line) if p.strip()]
            ids = [c for c in chain if re.fullmatch(r"[A-Za-z][A-Za-z0-9]*", c)]
            for a, b in zip(ids, ids[1:]):
                nodes.setdefault(a, a)
                nodes.setdefault(b, b)
                edges.append((a, b, elabel))
    return direction, nodes, edges


def _levels(nodes: dict[str, str], edges) -> list[list[str]]:
    """Xếp node thành từng tầng theo độ sâu — đủ cho sơ đồ ≤9 node."""
    depth = {n: 0 for n in nodes}
    for _ in range(len(nodes)):
        changed = False
        for a, b, _lbl in edges:
            if a in depth and b in depth and depth[b] < depth[a] + 1:
                depth[b] = depth[a] + 1
                changed = True
        if not changed:
            break
    out: list[list[str]] = []
    for n in nodes:
        d = depth[n]
        while len(out) <= d:
            out.append([])
        out[d].append(n)
    return [lv for lv in out if lv]


def _draw_diagram(slide, code: str, x, y, w, h) -> bool:
    """Vẽ sơ đồ bằng shape gốc PowerPoint — kéo và sửa được, không phải ảnh chết.

    Trả về False nếu không bóc được gì, để chỗ gọi còn xoay sang phương án khác.
    """
    direction, nodes, edges = parse_mermaid(code)
    if not nodes:
        return False
    levels = _levels(nodes, edges)
    if not levels:
        return False

    gap = 18
    at: dict[str, tuple[float, float, float, float]] = {}
    if direction == "LR":
        bw = (w - gap * (len(levels) - 1)) / len(levels)
        for i, lv in enumerate(levels):
            bh = (h - gap * (len(lv) - 1)) / len(lv)
            for j, n in enumerate(lv):
                at[n] = (x + i * (bw + gap), y + j * (bh + gap), bw, bh)
    else:
        bh = (h - gap * (len(levels) - 1)) / len(levels)
        for i, lv in enumerate(levels):
            bw = (w - gap * (len(lv) - 1)) / len(lv)
            for j, n in enumerate(lv):
                at[n] = (x + j * (bw + gap), y + i * (bh + gap), bw, bh)

    # cạnh vẽ trước để nằm dưới hộp
    for a, b, _lbl in edges:
        if a not in at or b not in at:
            continue
        ax, ay, aw, ah = at[a]
        bx, by, bw_, bh_ = at[b]
        if direction == "LR":
            p0, p1 = (ax + aw, ay + ah / 2), (bx, by + bh_ / 2)
        else:
            p0, p1 = (ax + aw / 2, ay + ah), (bx + bw_ / 2, by)
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        px(p0[0]), px(p0[1]), px(p1[0]), px(p1[1]))
        cn.line.color.rgb = MUTED
        cn.line.width = px(1.5)

    for n, label in nodes.items():
        if n not in at:
            continue
        bx, by, bw_, bh_ = at[n]
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    px(bx), px(by), px(bw_), px(bh_))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = LINE
        sh.line.width = px(1.5)
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = px(8)
        tf.margin_top = tf.margin_bottom = px(4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.2
        _rich_runs(p, label or n, size=17, color=INK)
    return True


# ------------------------------------------------------------- từng slide


def _add_notes(slide, text: str) -> None:
    if not (text or "").strip():
        return
    tf = slide.notes_slide.notes_text_frame
    tf.text = text.strip()


def _footer(slide, foot: str, no: str) -> None:
    _text(slide, PAD_X, H_PX - 34, W_PX - PAD_X * 2, 20, f"{foot} · {no}",
          size=13, color=MUTED, align=PP_ALIGN.CENTER)


def _rect(slide, x, y, w, h, fill, radius=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    if radius:
        try:                       # bo góc nhẹ thôi, mặc định của PowerPoint quá tròn
            sh.adjustments[0] = 0.08
        except Exception:          # noqa: BLE001
            pass
    return sh


def _hex(s: str) -> RGBColor:
    return RGBColor.from_string(s.lstrip("#").upper())


def _card(slide, c: dict, i: int, x, y, w, h) -> None:
    """Một thẻ pastel: chip màu, tiêu đề đậm, meta xám, các ý bên trong."""
    _rect(slide, x, y, w, h, _hex(theme.card_tint(i)))
    pad = 22
    # chip icon — python-pptx không nhúng SVG được, nên vẽ ô màu bo góc mang chữ
    # cái đầu của tên icon. Trên .pptx người dùng thay bằng icon riêng cũng dễ.
    chip = _rect(slide, x + pad, y + pad, 40, 40, _hex(theme.chip_color(i)))
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (c.get("icon") or "•")[:1].upper()
    r.font.size = px(20); r.font.bold = True
    r.font.name = FONT; r.font.color.rgb = WHITE

    tx = x + pad + 54
    tw = w - pad * 2 - 54
    title = (c.get("title") or "")
    meta = (c.get("meta") or "").strip()
    _text(slide, tx, y + pad + 4, tw, 34,
          title + (f"   {meta}" if meta else ""), size=21, bold=True)

    items = [b for b in (c.get("bullets") or []) if (b or "").strip()]
    if items:
        _bullets(slide, x + pad, y + pad + 54, w - pad * 2, h - pad * 2 - 54,
                 items, size=18, color=INK_2, gap=8)


def _render(slide, sl: dict, lay: str, ctx: dict, no: str) -> None:
    head = sl.get("headline") or ""
    sub = (sl.get("sub") or "").strip()
    cards = [c for c in (sl.get("cards") or []) if (c or {}).get("title")]
    bl = [b for b in (sl.get("bullets") or []) if (b or "").strip()]
    body_w = W_PX - PAD_X * 2

    if lay == "title":
        _text(slide, PAD_X, 214, body_w, 26,
              sl.get("eyebrow") or "BÁO CÁO SEMINAR", size=15, color=ACCENT, bold=True)
        _text(slide, PAD_X, 248, body_w, 130, head or ctx["title_vi"], size=52, bold=True)
        _text(slide, PAD_X, 386, body_w, 70, sub or ctx.get("title_en", ""),
              size=22, color=MUTED)
        _text(slide, PAD_X, 486, body_w, 30, ctx.get("venue", ""), size=18, color=INK_2)
        _text(slide, PAD_X, 514, body_w, 30, ctx.get("source", ""), size=18, color=MUTED)
        return

    if lay == "section":
        _text(slide, PAD_X, 276, body_w, 26, sl.get("eyebrow") or "PHẦN",
              size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, PAD_X, 310, body_w, 110, head, size=52, bold=True,
              align=PP_ALIGN.CENTER)
        if sub:
            _text(slide, PAD_X, 424, body_w, 44, sub, size=20, color=MUTED,
                  align=PP_ALIGN.CENTER)
        _footer(slide, ctx["foot"], no)
        return

    if lay == "closing":
        _text(slide, PAD_X, 268, body_w, 90, head, size=52, bold=True,
              align=PP_ALIGN.CENTER)
        if sub:
            _text(slide, PAD_X, 360, body_w, 40, sub, size=20, color=MUTED,
                  align=PP_ALIGN.CENTER)
        if bl:
            _bullets(slide, PAD_X, 412, body_w, 140, bl, size=20, color=INK_2,
                     marker="")
        _footer(slide, ctx["foot"], no)
        return

    # ---- đầu slide chung cho mọi bố cục nội dung
    y = PAD_TOP
    if (eb := (sl.get("eyebrow") or "").strip()):
        _text(slide, PAD_X, y, body_w, 24, eb, size=15, color=ACCENT, bold=True)
        y += 28
    _text(slide, PAD_X, y, body_w, 108, head, size=42, bold=True)
    y += 60 if len(head) < 58 else 106
    if sub:
        _text(slide, PAD_X, y, body_w, 40, sub, size=20, color=MUTED)
        y += 44
    y += 14
    avail = H_PX - PAD_BOT - y

    if lay == "agenda":
        rows = cards or [{"title": b} for b in bl]
        rh = min(84.0, avail / max(len(rows), 1))
        for i, c in enumerate(rows):
            ry = y + i * rh
            _rect(slide, PAD_X, ry, body_w, rh - 10, _hex(theme.card_tint(i)))
            n = _rect(slide, PAD_X + 16, ry + (rh - 10 - 46) / 2, 46, 46,
                      _hex(theme.chip_color(i)))
            ntf = n.text_frame
            ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
            np = ntf.paragraphs[0]; np.alignment = PP_ALIGN.CENTER
            nr = np.add_run(); nr.text = str(i + 1)
            nr.font.size = px(22); nr.font.bold = True
            nr.font.name = FONT; nr.font.color.rgb = WHITE
            desc = next((x for x in (c.get("bullets") or []) if (x or "").strip()), "")
            _text(slide, PAD_X + 78, ry + 10, body_w - 96, 30, c.get("title") or "",
                  size=21, bold=True)
            if desc:
                _text(slide, PAD_X + 78, ry + 40, body_w - 96, 26, desc,
                      size=17, color=MUTED)
        _footer(slide, ctx["foot"], no)
        return

    # ---- chỗ dành cho callout và chú thích thuật ngữ ở đáy slide
    co = sl.get("callout") or {}
    has_co = bool(co.get("title") or co.get("body"))
    terms = sl.get("terms") or []
    bot = 0
    if has_co:
        bot += 78
    if terms:
        bot += 20 + 22 * len(terms)
    avail = max(120, avail - bot)

    def draw_cards(cx, cy, cw, ch) -> None:
        if not cards:
            return
        n = min(len(cards), 4)
        gap = 20
        if n <= 2:
            w = (cw - gap * (n - 1)) / n
            for i, c in enumerate(cards[:n]):
                _card(slide, c, i, cx + i * (w + gap), cy, w, ch)
        else:
            w = (cw - gap * (n - 1)) / n
            for i, c in enumerate(cards[:n]):
                _card(slide, c, i, cx + i * (w + gap), cy, w, ch)

    def draw_visual(vx, vy, vw, vh) -> None:
        if (fig := sl.get("figure")) and (p := store.image_path(ctx["doc_id"], fig)):
            cap = (sl.get("figure_note") or "").strip()
            ih = vh - (34 if cap else 0)
            _picture(slide, str(p), vx, vy, vw, ih)
            if cap:
                _text(slide, vx, vy + ih + 6, vw, 30, cap, size=17, color=MUTED,
                      align=PP_ALIGN.CENTER)
            return
        if (dia := (sl.get("diagram") or "").strip()):
            if _draw_diagram(slide, dia, vx, vy, vw, vh):
                return
        if (eq := (sl.get("equation") or "").strip()):
            _text(slide, vx, vy + vh / 2 - 24, vw, 48, eq, size=26,
                  align=PP_ALIGN.CENTER, font=MONO)

    if lay in ("figside", "split"):
        gap = 36
        tw = (body_w - gap) * 0.44
        vx = PAD_X + tw + gap
        if cards:
            draw_cards(PAD_X, y, tw, avail)
        elif bl:
            _bullets(slide, PAD_X, y, tw, avail, bl, size=20,
                     anchor=MSO_ANCHOR.MIDDLE)
        draw_visual(vx, y, body_w - tw - gap, avail)
    elif lay == "figwide":
        th = 0
        if cards:
            th = min(190.0, avail * 0.42)
            draw_cards(PAD_X, y, body_w, th)
        elif bl:
            th = min(150.0, avail * 0.38)
            _bullets(slide, PAD_X, y, body_w, th, bl, size=20)
        draw_visual(PAD_X, y + th + (16 if th else 0), body_w, avail - th - (16 if th else 0))
    elif lay == "figfull":
        draw_visual(PAD_X, y, body_w, avail)
    elif lay == "cards":
        draw_cards(PAD_X, y, body_w, min(avail, 300.0))
    else:
        _bullets(slide, PAD_X, y, body_w, avail, bl, size=21)

    # số liệu lớn
    st = [s for s in (sl.get("stats") or []) if (s or {}).get("value")][:2]
    if st and lay in ("cards", "list"):
        sy = y + min(avail, 300.0) + 16
        for i, s in enumerate(st):
            _text(slide, PAD_X + i * 320, sy, 300, 48, s.get("value") or "",
                  size=40, bold=True, color=ACCENT)
            _text(slide, PAD_X + i * 320, sy + 48, 300, 34, s.get("label") or "",
                  size=16, color=MUTED)

    yb = H_PX - PAD_BOT - bot
    if has_co:
        _rect(slide, PAD_X, yb, body_w, 66, _hex("#e0e9fd"))
        _text(slide, PAD_X + 22, yb + 12, body_w - 44, 26, co.get("title") or "",
              size=21, bold=True)
        if co.get("body"):
            _text(slide, PAD_X + 22, yb + 38, body_w - 44, 24, co["body"],
                  size=17, color=MUTED)
        yb += 78
    for t in terms:
        _text(slide, PAD_X, yb, body_w, 22, f"{t['en']} — {t['gloss']}",
              size=15, color=MUTED)
        yb += 22

    _footer(slide, ctx["foot"], no)


def _visual(slide, sl: dict, x, y, w, h, ctx: dict) -> None:
    """Hình cắt từ bài, sơ đồ vẽ lại bằng shape, hoặc công thức."""
    if (fig := sl.get("figure")):
        p = store.image_path(ctx["doc_id"], fig)
        if p is not None:
            _picture(slide, str(p), x, y, w, h)
            return
    if (dia := (sl.get("diagram") or "").strip()):
        if _draw_diagram(slide, dia, x, y, w, h):
            return
    if (eq := (sl.get("equation") or "").strip()):
        _text(slide, x, y + h / 2 - 30, w, 60, eq, size=30,
              align=PP_ALIGN.CENTER, font=MONO)


def _picture(slide, path: str, x, y, w, h) -> None:
    """Chèn ảnh, giữ đúng tỉ lệ và căn giữa trong khung cho trước."""
    from PIL import Image  # đi kèm pymupdf/docling, không thêm phụ thuộc mới
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:  # noqa: BLE001
        iw, ih = 4, 3
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    slide.shapes.add_picture(path, px(x + (w - dw) / 2), px(y + (h - dh) / 2),
                             px(dw), px(dh))


# ------------------------------------------------------------------ vào ra


def build(doc: dict) -> bytes:
    """Cả bộ slide thành một file .pptx."""
    from . import pipeline

    brief = doc.get("brief") or {}
    slides = doc.get("slides") or {}
    deck = list(slides.get("deck") or [])
    backup = list(slides.get("backup") or [])

    prs = Presentation()
    # `Pt` trả về Length tính bằng EMU sẵn — đừng nhân 12700 lần nữa
    prs.slide_width = px(W_PX)
    prs.slide_height = px(H_PX)
    blank = prs.slide_layouts[6]        # bố cục trống — ta tự đặt mọi khung chữ

    title_vi = brief.get("title_vi") or doc.get("title") or "Bài báo"
    ctx = {
        "doc_id": doc["id"],
        "title_vi": title_vi,
        "title_en": doc.get("title") or "",
        "venue": (brief.get("venue_guess") or "")[:90],
        "source": (doc.get("source") or "")[:95],
        "foot": title_vi[:70],
    }

    secs = [s for s in deck if (s.get("kind") or "") == "section"]
    for n, s in enumerate(secs, 1):
        s["eyebrow"] = f"PHẦN {n} / {len(secs)}"
    # thuật ngữ lần đầu xuất hiện thì gắn chú thích, lấy từ bảng đã chốt ở brief
    pipeline.attach_terms(doc, deck)
    pipeline.attach_terms(doc, backup)

    for group, tag in ((deck, ""), (backup, "D")):
        for i, sl in enumerate(group, 1):
            slide = prs.slides.add_slide(blank)
            _render(slide, sl, pipeline.slide_layout(sl, doc["id"]), ctx, f"{tag}{i}")
            _add_notes(slide, sl.get("notes") or "")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
